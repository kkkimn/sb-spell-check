import re
import json
import time
import difflib
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openai import OpenAI
import fitz  # PyMuPDF


# ==========================================
# 공통 헬퍼: 그룹 도형 재귀 순회
# ==========================================
def iter_shapes(shapes):
    """
    그룹 도형(GROUP) 안에 중첩된 shape까지 재귀적으로 모두 yield 한다.
    PPT는 그룹화가 매우 흔해서 단순 for shape in slide.shapes 만 돌면
    그룹 안의 텍스트가 30~50%까지 누락될 수 있다.
    """
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def _safe_shape_text(shape):
    """
    어떤 종류의 도형이든 텍스트를 안전하게 읽어온다.
    has_text_frame 체크가 실패하는 특수 도형(placeholder, autoshape 등)도 처리.
    텍스트가 없거나 읽기 실패 시 빈 문자열 반환.
    """
    # 1차: text_frame 직접 접근
    try:
        if shape.has_text_frame:
            t = shape.text_frame.text
            if t:
                return t
    except Exception:
        pass
    # 2차: shape.text 속성 직접 시도
    try:
        t = shape.text
        if t:
            return t
    except Exception:
        pass
    return ""


def _safe_is_bottom(shape, slide_height):
    """shape.top이 None인 경우를 안전하게 처리"""
    top = getattr(shape, "top", None)
    if top is None or slide_height is None:
        return False
    return top > slide_height * 0.6


def _normalize(text):
    """
    공백·줄바꿈·탭을 모두 제거한 비교용 문자열.
    PPT의 세로쓰기 셀(예: '내\\n레\\n이\\n션')을 정상 매칭하기 위해 사용.
    """
    if not text:
        return ""
    return re.sub(r'\s+', '', text)


def _detect_speaker(text):
    """
    텍스트에서 화자 키워드를 찾아 화자명을 반환한다.
    공백·줄바꿈을 제거한 후 매칭하므로 세로쓰기 셀에도 동작한다.
    '내레이션'/'나레이션'/'narration' 은 '기타'로 매핑.
    매칭 실패 시 None.
    """
    if not text:
        return None
    norm = _normalize(text)
    if not norm:
        return None
    # 우선순위: 명시적 화자 > 내레이션
    if "교수" in norm:
        return "교수"
    if "선생님" in norm:
        return "선생님"
    if "성우" in norm:
        return "성우"
    if "내레이션" in norm or "나레이션" in norm or "narration" in norm.lower():
        return "기타"
    return None


def _is_narration_label(text):
    """셀/라벨 텍스트가 '내레이션' 류 라벨인지 판단 (세로쓰기 대응)."""
    if not text:
        return False
    norm = _normalize(text).lower()
    return "내레이션" in norm or "나레이션" in norm or "narration" in norm


def _strip_numbering_markers(text):
    """
    본문에서 '#1', '#2', '# 3', '#10.', '#1)', '#1:' 같은 번호 마커를 모두 제거하고,
    그 결과로 생긴 다중 공백을 정리한다.
    """
    if not text:
        return text
    # # 다음에 (선택적 공백) 숫자 (선택적 종결 부호) 패턴 제거
    text = re.sub(r'#\s*\d+\s*[).:．。]?', '', text)
    # 다중 공백 정리
    text = re.sub(r'[ \t]+', ' ', text)
    # 빈 줄 정리 (마커 제거 후 줄바꿈만 남는 경우 대비)
    text = re.sub(r'\n\s*\n', '\n', text)
    return text.strip()


def _clean_punctuation(text):
    """
    맞춤법 교정 결과에서 흔히 발생하는 구두점 중복/오류를 정리한다.
    
    GPT가 원문 '~다'에 마침표를 추가해 '~다.'로 반환하면, 원문에 이미 마침표가
    있을 경우 치환 결과가 '~다..' 같은 이중 마침표가 되는 경우가 있다.
    
    처리 항목:
    - '..' → '.'  (단, '...' 같은 말줄임표 3개 이상은 보존)
    - ',,', '!!', '??' 도 동일하게 단일화
    - 마침표/콤마/물음표/느낌표 앞의 불필요한 공백 제거 (' .' → '.')
    - 다중 공백 정리
    """
    if not text:
        return text
    # 이중 마침표 → 단일 (앞뒤가 점이 아닌 경우만 = 말줄임표 보존)
    text = re.sub(r'(?<!\.)\.\.(?!\.)', '.', text)
    # 이중 콤마/느낌표/물음표
    text = re.sub(r',{2,}', ',', text)
    text = re.sub(r'!{2,}', '!', text)
    text = re.sub(r'\?{2,}', '?', text)
    # 마침표/콤마/느낌표/물음표/콜론 앞의 공백 제거
    text = re.sub(r' +([.,!?:])', r'\1', text)
    # 다중 공백 정리
    text = re.sub(r' {2,}', ' ', text)
    return text


def apply_corrections_to_narration(narration_dict, corrections_dict):
    """
    추출된 원본 대본 dict에 맞춤법 교정을 적용하고 구두점을 정리한다.
    PPT/PDF 공통 사용.
    
    파라미터:
        narration_dict: {화자: [대본 라인 리스트]} — extract_narrations() 결과
        corrections_dict: {원문: 교정문} — get_openai_corrections_*() 결과
    
    반환: 교정 + 구두점 정리가 적용된 대본 dict
    """
    if not corrections_dict:
        # 교정 사항이 없어도 구두점 정리는 수행
        result = {}
        for speaker, lines in narration_dict.items():
            result[speaker] = [_clean_punctuation(line) for line in lines]
        return result
    
    # 긴 키부터 정렬해서 짧은 단어가 다른 단어를 망가뜨리지 않도록
    sorted_items = sorted(corrections_dict.items(), key=lambda x: -len(x[0]))
    
    result = {}
    for speaker, lines in narration_dict.items():
        new_lines = []
        for line in lines:
            new_line = line
            for old_txt, new_txt in sorted_items:
                if old_txt and old_txt in new_line:
                    new_line = new_line.replace(old_txt, new_txt)
            new_line = _clean_punctuation(new_line)
            new_lines.append(new_line)
        result[speaker] = new_lines
    
    return result


def _parse_speaker_label(content):
    """
    내레이션 본문 앞부분에서 화자 라벨을 찾아낸다.
    
    인식 규칙: 본문의 맨 앞이 'XXX :' 또는 'XXX:' 형태이면 XXX를 라벨로 본다.
      - XXX는 공백·콜론·줄바꿈을 포함하지 않는 1~10자
      - 구분자는 ':' 또는 '：' (전각 콜론)
      - 구분자 앞뒤의 공백은 허용
    
    인식되면 (label, 나머지_내용)을, 인식 실패 시 (None, 원본_내용)을 반환한다.
    """
    if not content:
        return None, content
    match = re.match(r'^\s*([^\s:：\n]{1,10})\s*[:：]\s*(.+)', content, flags=re.DOTALL)
    if match:
        label = match.group(1).strip()
        rest = match.group(2).strip()
        if label and rest:
            return label, rest
    return None, content.strip()


# 다중 화자 분리용 - 중간에서 인식할 화자 키워드 whitelist
# 맨 앞(문장 시작)의 라벨은 permissive하게 아무 1-10자를 허용하지만,
# 중간의 라벨은 이 whitelist 키워드로 끝나는 것만 인정한다.
# → '다음 내용:', '시간:', '목차:' 같은 일반 콜론을 화자로 오인하는 것 방지
_SPEAKER_KEYWORDS = [
    '교수', '선생님', '선생', '성우', '강사', '학생',
    '진행자', '사회자', '아나운서', '내레이터', '해설자', '해설',
    '박사', '강연자', '발표자',
]

_SPEAKER_KEYWORDS_ALT = '|'.join(sorted(_SPEAKER_KEYWORDS, key=len, reverse=True))

# 맨 앞 permissive 매칭에서 '화자가 아님이 명백한' 일반 명사 블랙리스트
# (이 단어들이 맨 앞 'XX:' 형태로 나타나면 라벨로 인정하지 않음)
_NON_SPEAKER_LABELS = frozenset([
    '시간', '장소', '일시', '일자', '날짜', '기간', '대상',
    '목차', '순서', '내용', '주제', '제목', '저자', '작성자', '편집자',
    '비고', '주의', '경고', '참조', '참고', '출처', '예시', '예제', '정답', '문제',
    '첫째', '둘째', '셋째', '넷째', '다섯째',
    '가격', '금액', '비용', '합계', '총계',
    '주소', '연락처', '전화', '이메일', '메일', '팩스',
    'http', 'https', 'www', 'url', 'email', 'mail', 'tel', 'phone', 'fax', 'note',
])

# 맨 앞 라벨 (permissive): 1~10자 아무 비공백 문자 + 콜론
_LEADING_LABEL_RE = re.compile(r'^\s*([^\s:：\n]{1,10})\s*[:：]\s*')

# 중간 라벨 (whitelist): [한글 0-3자 접두사] + 화자 키워드 + [님?] + 콜론
# - 이름 접두사 지원: '김교수:', '박선생님:'
# - 호칭 접미사 지원: '교수님:'
_MIDDLE_LABEL_RE = re.compile(
    rf'([가-힣]{{0,3}}(?:{_SPEAKER_KEYWORDS_ALT}))님?\s*[:：]\s*'
)


def _split_multi_speaker(content):
    """
    본문 안에서 화자 라벨을 찾아 세그먼트 리스트로 분리한다.
    
    동작:
    - 맨 앞의 라벨은 permissive (1-10자 아무거나 + 콜론)
    - 중간의 라벨은 whitelist (교수/선생님/강사/성우/... + 콜론)
      → 이름 접두사('김교수:'), 호칭 접미사('교수님:') 모두 지원
    
    반환: [(label_or_None, text), ...]
    
    예)
      "교수: 안녕하세요" 
        → [('교수', '안녕하세요')]
      "교수: 안녕하세요 선생님: 반갑습니다" 
        → [('교수', '안녕하세요'), ('선생님', '반갑습니다')]
      "안녕하세요 교수: 시작합시다" 
        → [(None, '안녕하세요'), ('교수', '시작합시다')]
      "오늘은 이것을 배웁니다" 
        → [(None, '오늘은 이것을 배웁니다')]
      "다음 내용: 첫째, 둘째, 셋째" 
        → [(None, '다음 내용: 첫째, 둘째, 셋째')]  # 위양성 방지
      "김교수: A 박선생님: B"
        → [('김교수', 'A'), ('박선생님', 'B')]
    """
    if not content:
        return []
    content = content.strip()
    if not content:
        return []
    
    # 1) 맨 앞 permissive 라벨 시도 (단, 블랙리스트에 있으면 거부)
    leading_label = None
    body_start = 0
    m = _LEADING_LABEL_RE.match(content)
    if m:
        candidate = m.group(1).strip()
        if candidate and candidate.lower() not in _NON_SPEAKER_LABELS:
            leading_label = candidate
            body_start = m.end()
    
    remaining = content[body_start:]
    
    # 2) 나머지에서 whitelist 기반 split point 찾기
    split_points = []
    for match in _MIDDLE_LABEL_RE.finditer(remaining):
        # 문자열 시작 또는 공백/줄바꿈 뒤에서만 유효 (단어 중간 매칭 방지)
        if match.start() == 0 or remaining[match.start() - 1] in ' \t\n':
            split_points.append(match)
    
    # 3) 세그먼트 구성
    if not split_points:
        body = remaining.strip()
        if not body:
            return []
        return [(leading_label, body)]
    
    segments = []
    
    # 첫 split point 이전 텍스트 (leading_label 하에 소속)
    first_chunk = remaining[:split_points[0].start()].strip()
    if first_chunk:
        segments.append((leading_label, first_chunk))
    
    # 각 split point 이후의 세그먼트
    for i, sp in enumerate(split_points):
        label = sp.group(1).strip()
        chunk_start = sp.end()
        chunk_end = split_points[i + 1].start() if i + 1 < len(split_points) else len(remaining)
        chunk = remaining[chunk_start:chunk_end].strip()
        if chunk:
            segments.append((label, chunk))
    
    return segments


# 내레이션이 아닌 노이즈 텍스트를 걸러내기 위한 키워드 (대소문자 무시)
_NON_NARRATION_KEYWORDS = (
    "버전", "version", "ver.", "v1.", "v2.", "v3.",
    "수정일", "작성일", "갱신일", "제작일", "제작자", "작성자",
    "화면설명", "화면 설명", "[화면]", "화면 ]",
    "차시", "단원", "학습목표", "학습 목표",
    "출처", "참고문헌", "참고 문헌",
    "페이지", "page", "p.",
    "copyright", "©", "all rights reserved",
    # placeholder 텍스트 — 비어있는 템플릿 자리 표시
    "내용을 입력", "텍스트를 입력", "여기를 클릭", "제목을 입력",
    "내용 입력", "텍스트 입력", "이미지 번호", "이미지번호",
    "click here", "enter text",
)


def _is_non_narration_noise(text):
    """
    (슬라이드 노트 필터링용) 버전 정보/화면 설명/페이지 푸터 등 
    내레이션이 아닌 것이 명백한 텍스트인지 판단. 길이도 함께 체크.
    """
    if not text:
        return True
    lower = text.lower()
    for kw in _NON_NARRATION_KEYWORDS:
        if kw in lower:
            return True
    if len(text.strip()) < 10:
        return True
    return False


def _cell_is_noise(text):
    """
    (셀/도형 처리용) 길이 체크 없이 노이즈 키워드만 검사.
    짧은 화자 셀('교수', '선생님')이 잘못 제외되지 않도록.
    """
    if not text:
        return True
    lower = text.lower()
    for kw in _NON_NARRATION_KEYWORDS:
        if kw in lower:
            return True
    return False


def _sorted_corrections(corrections_dict):
    """
    교정 딕셔너리를 '키 길이 내림차순'으로 정렬해서 반환.
    이렇게 해야 '이→의' 같은 짧은 치환이 '이것은' 안의 '이'를 망가뜨리지 않는다.
    """
    return sorted(corrections_dict.items(), key=lambda x: -len(x[0]))


# ==========================================
# 화자별 대본 추출 (PPT)
# ==========================================
def extract_narrations(prs):
    """
    PPT에서 '내레이션' 영역의 대본을 추출한다.
    
    세 가지 전략을 순차 시도한다:
    1. 진짜 표(table) 안에 '내레이션' 라벨 셀이 있는 경우
    2. '내레이션' 라벨이 텍스트 박스로 되어 있고 본문도 별개 텍스트 박스인 '가짜 표' 레이아웃
       (슬라이드 본체 + 슬라이드 레이아웃 + 마스터까지 검색)
    3. (1, 2 모두 실패) 휴리스틱 폴백 — 슬라이드 하단의 충분히 길고 넓은 텍스트 도형을
       내레이션 본문으로 간주. '내레이션' 라벨이 개별 글자 박스로 분리되어 있어서
       라벨 검출이 불가능한 템플릿용.
    
    본문 앞부분의 화자 라벨('교수:', '선생님:', '강사:' 등)을 동적으로 그룹핑한다.
    라벨이 없는 본문은 '라벨없음' 그룹으로 모은다.
    """
    narrations = {}  # dict[label -> list[entry]]  (동적으로 key가 추가됨)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        # 전략 1: 진짜 표
        found = _extract_narration_from_tables(slide, slide_num, narrations)
        # 전략 2: 텍스트 박스 레이블 + 공간 매칭
        if not found:
            found = _extract_narration_from_text_frames(slide, slide_num, narrations)
        # 전략 3: 휴리스틱 폴백 (라벨 없이 하단 본문 직접 추출)
        if not found:
            _extract_narration_from_bottom_strip(
                slide, slide_num, narrations, slide_width, slide_height
            )
    
    return narrations


def _contains_speaker_keyword(text):
    """
    텍스트에 whitelist 화자 키워드(교수/선생님/강사/...)가 하나라도 포함되어 있는지 검사.
    공백·줄바꿈을 제거한 뒤 검사하므로 세로쓰기 셀에도 동작.
    """
    if not text:
        return False
    norm = _normalize(text)
    for kw in _SPEAKER_KEYWORDS:
        if kw in norm:
            return True
    return False


def _add_narration_entry(narrations, slide_num, content_text, speaker_override=None):
    """
    공통: 본문 텍스트를 정리하고 화자별 그룹에 추가.
    한 덩어리 안에 여러 화자가 섞여있으면 _split_multi_speaker가 자동 분리.
    
    speaker_override: 표의 별도 셀에서 감지한 화자 텍스트 (예: "김영곤 선생님").
                      본문 자체에 inline 화자 라벨이 없을 때 이 값을 그룹 key로 사용한다.
                      본문에 inline 라벨이 있으면 그쪽이 우선.
    """
    content_text = _strip_numbering_markers(content_text)
    if not content_text:
        return
    
    segments = _split_multi_speaker(content_text)
    for label, body in segments:
        if not body:
            continue
        # 본문 자체에 inline 라벨이 있으면 그것을 우선, 없으면 speaker_override
        if label:
            group_key = label
        elif speaker_override:
            group_key = speaker_override
        else:
            group_key = "라벨없음"
        narrations.setdefault(group_key, []).append(
            f"[슬라이드 {slide_num}] {group_key} :\n{body}"
        )


def _process_cell_group(cell_texts, slide_num, narrations, source_tag=""):
    """
    여러 텍스트 조각(셀 또는 도형 텍스트)을 받아서 화자/본문 분리 후 narrations에 추가.
    전략 1, 2, 3 모두에서 공통으로 사용.
    
    로직:
    1. 노이즈 키워드 텍스트는 제외 (버전, 화면설명, placeholder 등)
    2. 짧고(≤20자) 화자 키워드 포함 → 화자 셀 후보 (첫 번째만 사용)
    3. 그 외 → 본문 후보. 가장 긴 것을 본문으로 선택.
    4. 본문이 없고 화자 셀만 있으면 화자 셀을 본문으로 사용
    
    반환: 성공 여부.
    """
    speaker_cell_text = None
    body_candidates = []
    filtered_out = []  # 진단용
    
    for t in cell_texts:
        if not t:
            continue
        t = t.strip()
        if not t:
            continue
        if _cell_is_noise(t):
            filtered_out.append(t)
            continue
        if len(t) <= 20 and _contains_speaker_keyword(t):
            if speaker_cell_text is None:
                speaker_cell_text = t
                continue
        body_candidates.append((len(t), t))
    
    if body_candidates:
        body_candidates.sort(key=lambda x: -x[0])
        body_text = body_candidates[0][1]
    elif speaker_cell_text:
        body_text = speaker_cell_text
        speaker_cell_text = None
    else:
        return False
    
    # 진단: 어떤 본문이 선택됐는지 출력 (디버깅용)
    body_preview = body_text[:50] + ("..." if len(body_text) > 50 else "")
    speaker_info = f" 화자='{speaker_cell_text}'" if speaker_cell_text else ""
    print(f"   [추출] 슬라이드 {slide_num} {source_tag}:{speaker_info} 본문='{body_preview}'")
    
    _add_narration_entry(
        narrations, slide_num, body_text,
        speaker_override=speaker_cell_text
    )
    return True


def _extract_narration_from_tables(slide, slide_num, narrations):
    """
    전략 1: 진짜 표 안에서 '내레이션' 라벨 행 찾기.
    라벨 셀을 제외한 나머지 셀들을 _process_cell_group에 넘겨 화자/본문 분리.
    """
    found_any = False
    for shape in iter_shapes(slide.shapes):
        if not shape.has_table:
            continue
        
        for row in shape.table.rows:
            cells = list(row.cells)
            if len(cells) < 2:
                continue
            
            # 행 안에서 '내레이션' 라벨이 있는 셀 찾기 (위치 무관)
            narration_label_idx = None
            for idx, cell in enumerate(cells):
                if _is_narration_label(cell.text):
                    narration_label_idx = idx
                    break
            
            if narration_label_idx is None:
                continue
            
            # 라벨 셀을 제외한 나머지 셀들을 공통 헬퍼로 처리
            cell_texts = [
                cell.text for idx, cell in enumerate(cells) 
                if idx != narration_label_idx
            ]
            if _process_cell_group(cell_texts, slide_num, narrations, source_tag="표"):
                found_any = True
    
    return found_any


def _extract_narration_from_text_frames(slide, slide_num, narrations):
    """
    전략 2 (폴백): 텍스트 박스로 만든 '가짜 표' 레이아웃 처리.
    
    동작 방식:
    - 슬라이드 + 슬라이드 레이아웃 + 슬라이드 마스터 모두 스캔
      (한국 교육용 템플릿은 '내레이션' 라벨 같은 고정 요소를 레이아웃/마스터에 두는 경우가 많음)
    - '내레이션' 라벨 도형과 일반 텍스트 도형을 분리
    - 각 라벨 도형에 대해, 같은 가로 행에 위치하면서 라벨의 오른쪽으로 뻗어있는
      텍스트 도형들을 모아 본문으로 사용
    
    추출 성공 여부 반환.
    """
    label_shapes = []
    content_candidates = []
    
    # 슬라이드 본체에서 도형 수집
    for shape in iter_shapes(slide.shapes):
        text = _safe_shape_text(shape).strip()
        if not text:
            continue
        if _is_narration_label(text):
            label_shapes.append(shape)
        else:
            content_candidates.append(shape)
    
    # 라벨이 슬라이드 본체에서 안 보이면, 레이아웃과 마스터에서도 찾아본다
    if not label_shapes:
        layout_sources = []
        try:
            layout_sources.append(slide.slide_layout)
        except Exception:
            pass
        try:
            layout_sources.append(slide.slide_layout.slide_master)
        except Exception:
            pass
        
        for src in layout_sources:
            try:
                for shape in iter_shapes(src.shapes):
                    text = _safe_shape_text(shape).strip()
                    if text and _is_narration_label(text):
                        label_shapes.append(shape)
            except Exception:
                continue
    
    if not label_shapes:
        return False
    
    matched_any = False
    for label_shape in label_shapes:
        label_left = getattr(label_shape, 'left', None)
        label_top = getattr(label_shape, 'top', None)
        label_width = getattr(label_shape, 'width', None)
        label_height = getattr(label_shape, 'height', None)
        
        if None in (label_left, label_top, label_width, label_height):
            continue
        if label_height <= 0 or label_width <= 0:
            continue
        
        label_bottom = label_top + label_height
        
        row_items = []
        for cand in content_candidates:
            c_left = getattr(cand, 'left', None)
            c_top = getattr(cand, 'top', None)
            c_width = getattr(cand, 'width', None)
            c_height = getattr(cand, 'height', None)
            
            if None in (c_left, c_top, c_width, c_height):
                continue
            
            c_center_y = c_top + c_height / 2
            if not (label_top <= c_center_y <= label_bottom):
                continue
            
            if c_left + c_width <= label_left:
                continue
            
            row_items.append((c_left, cand))
        
        if not row_items:
            continue
        
        row_items.sort(key=lambda x: x[0])
        # 합치지 말고 각 박스 텍스트를 리스트로 수집 → 공통 헬퍼가 화자/본문 자동 분리
        box_texts = [_safe_shape_text(c).strip() for _, c in row_items]
        box_texts = [t for t in box_texts if t]
        if not box_texts:
            continue
        
        if _process_cell_group(box_texts, slide_num, narrations, source_tag="가짜표"):
            matched_any = True
    
    return matched_any


def _looks_like_narration_body(text):
    """
    텍스트가 '내레이션 본문'처럼 보이는지 휴리스틱 검사.
    #N 번호 마커가 있거나, 한국어 서술형 종결어미가 있으면 True.
    (화면설명 메모, 버전 정보, URL 등과 구분하기 위함)
    """
    if not text or len(text) < 15:
        return False
    # #N 마커 (내레이션 본문에서 매우 흔함)
    if re.search(r'#\s*\d+', text):
        return True
    # 서술형 종결어미 (문장이 끝나는 패턴)
    if re.search(r'(습니다|입니다|합니다|됩니다|였습니다|했습니다|겠습니다)[.。]?', text):
        return True
    return False


def _extract_narration_from_bottom_strip(slide, slide_num, narrations, slide_width, slide_height):
    """
    전략 3 (휴리스틱 폴백): '내레이션' 라벨을 못 찾아도 슬라이드 하단의
    텍스트 박스들에서 내레이션을 추출한다.
    
    동작 방식:
    - 슬라이드 하단 35% 영역의 모든 텍스트 도형 수집
    - 그 중 '내레이션 본문처럼 보이는' 텍스트(#N 마커 또는 서술형 종결어미)가
      하나 이상 있어야만 작동 (false positive 방지)
    - 수집한 텍스트들을 _process_cell_group으로 넘겨서 화자 셀/본문 셀 분리
    
    → 전략 1과 2가 모두 실패한 경우에도, 하단 영역에 '김영곤 선생님 | 본문' 같은
      구조가 있으면 화자별 그룹핑까지 해준다.
    """
    if not slide_width or not slide_height:
        print(f"   [진단] 슬라이드 {slide_num}: 슬라이드 크기를 알 수 없어 하단 영역 검색 불가")
        return False
    
    bottom_threshold = slide_height * 0.65
    
    # 하단 영역의 모든 텍스트 박스 수집
    bottom_texts = []
    has_narration_signal = False
    
    for shape in iter_shapes(slide.shapes):
        text = _safe_shape_text(shape).strip()
        if not text:
            continue
        
        top = getattr(shape, 'top', None)
        left = getattr(shape, 'left', None)
        width = getattr(shape, 'width', None)
        if top is None or left is None or width is None:
            continue
        if top < bottom_threshold:
            continue
        
        # 오른쪽 사이드바(화면설명 등) 제외: 왼쪽 위치가 슬라이드 폭의 70%를 넘으면 스킵
        if left > slide_width * 0.7:
            continue
        
        # 노이즈 키워드 제외 (화면설명, 버전, 저작권 등)
        if _cell_is_noise(text):
            continue
        
        bottom_texts.append(text)
        
        # 이 텍스트가 내레이션 본문처럼 보이는지 체크
        if _looks_like_narration_body(text):
            has_narration_signal = True
    
    # 내레이션 본문 신호가 없으면 추출 중단 (false positive 방지)
    if not has_narration_signal:
        if bottom_texts:
            print(f"   [진단] 슬라이드 {slide_num}: 하단에 텍스트 {len(bottom_texts)}개 있으나 "
                  f"내레이션 본문 신호(#N 마커, 서술형 종결어미) 없음 → 추출 중단")
        else:
            print(f"   [진단] 슬라이드 {slide_num}: 모든 추출 전략 실패. 하단 영역 텍스트 없음.")
        return False
    
    # 공통 헬퍼로 화자 셀/본문 셀 분리 → 그룹핑
    return _process_cell_group(
        bottom_texts, slide_num, narrations, source_tag="하단폴백"
    )


# ==========================================
# 공통: 한국어 맞춤법 교정 시스템 프롬프트 빌더
# ==========================================
def _build_system_prompt(custom_dict, doc_kind="슬라이드"):
    """
    한국어 맞춤법 교정 프롬프트를 생성한다.
    회수율(recall)을 최우선으로 한다 — 의심되는 모든 오류를 적극적으로 잡아낸다.
    """
    custom_dict_prompt = ""
    if custom_dict and len(custom_dict) > 0:
        custom_dict_prompt = (
            "\n\n[예외 처리 — 사용자 맞춤법 사전]\n"
            "다음 단어들은 사용자의 의도적인 고유명사·예외 단어이다. "
            "원문에 등장하면 절대 띄어쓰기·맞춤법을 수정하지 말고 원형 그대로 보존해라. "
            "이 단어들 외의 모든 오류는 적극적으로 잡아라:\n"
            + ", ".join(custom_dict)
        )

    return f"""너는 한국어 맞춤법 교정 전문가야. {doc_kind} 텍스트를 **철저히 정밀 검사**해서 모든 오류를 빠짐없이 잡아내라. 작은 오류 하나도 놓치면 안 된다.

[검사 대상 — 다음을 모두 포함하되 여기에만 한정되지 않는다]
1. 맞춤법: 되/돼, 안/않, 왠/웬, 데/대, 던/든, ~로서/~로써, ~므로/~음으로, 잇/있, 햇/했, 것/거 등
2. 띄어쓰기: 조사 붙여쓰기, 의존명사 띄어쓰기, 보조용언, 단위명사, 합성어 등
3. 조사·어미: 을/를, 이/가, 은/는, ~에/~에서, ~으로/~로, 종결어미 활용
4. 외래어 표기법: 컨텐츠→콘텐츠, 악세사리→액세서리, 메세지→메시지, 후라이→프라이, 까페→카페, 초콜렛→초콜릿, 케잌→케이크, 화이팅→파이팅 등
5. 자판 오타·탈자: 하겟습니다→하겠습니다, 됬다→됐다, 갔다→갖다(문맥) 등
6. 사이시옷: 나뭇가지, 햇볕, 등굣길 등
7. 한자어 오용 및 부자연스러운 어순
8. 비표준어·구어체 오류: 어떻해→어떡해, 이쁘다→예쁘다(문맥), 짜장면/자장면 등

[적극성 원칙 — 매우 중요]
- **조금이라도 어색하거나 의심스러우면 일단 교정안을 제시해라.** 사용자가 최종 검토하므로 거짓 양성(false positive)이 거짓 음성(false negative)보다 훨씬 낫다.
- 한 페이지에서 발견되는 오류 개수에 인위적인 상한을 두지 마라. 20개든 50개든 발견된 모든 오류를 반환해라.
- 같은 유형의 오류가 반복되면 모두 다 잡아라. "비슷한 거 이미 잡았으니까" 라고 생각하고 건너뛰지 마라.

[출력 형식]
- 반드시 '순수 JSON 객체'로 응답. 교정할 것이 없으면 {{}} 만 반환.
- key는 원문, value는 교정문. **가능하면 짧은 어절 단위(1~3어절)**로 잘라서 반환해라. 단, 띄어쓰기 오류처럼 더 긴 단위가 자연스러우면 그렇게 해도 된다.
- key와 value가 완전히 같으면 결과에서 빼라.
- key는 반드시 원문에 **글자 그대로** 존재해야 한다. 임의로 공백을 넣거나 빼지 마라.

[하지 말아야 할 것 — 최소한의 가드레일]
- 화면 설명/UI 라벨/슬라이드 제목 같은 개조식·명사형 문장을 억지로 '~합니다' 완성형으로 바꾸지 마라. 단, 그 안의 맞춤법·띄어쓰기·외래어 오류는 반드시 잡아라.
- 의미를 바꾸는 윤문(rewriting)은 하지 마라. 오로지 표기 오류만 고쳐라.

[예시]
입력: "이 자료는 챗지피티를 활용 하여 만들어 졋습니다. 컨텐츠 가 풍부 합니다"
출력: {{"활용 하여": "활용하여", "만들어 졋습니다": "만들어졌습니다", "컨텐츠 가": "콘텐츠가", "풍부 합니다": "풍부합니다"}}

입력: "다음주 회의때 발표 할께요. 자료는 메일로 보내드릴께요"
출력: {{"다음주": "다음 주", "회의때": "회의 때", "발표 할께요": "발표할게요", "보내드릴께요": "보내드릴게요"}}

입력: "어떻해 해야 할 지 모르겟어요. 왠지 느낌이 안 좋네요"
출력: {{"어떻해": "어떡해", "할 지": "할지", "모르겟어요": "모르겠어요"}}

입력: "악세사리 코너에서 초콜렛 케잌을 샀어요"
출력: {{"악세사리": "액세서리", "초콜렛": "초콜릿", "케잌을": "케이크를"}}{custom_dict_prompt}
"""


def _is_custom_dict_violation(k, v, custom_dict):
    """
    사용자 사전 단어가 원문에 등장했는데 교정 후에는 보존되지 않았으면 위반.
    
    검사 방식:
    1. 단어가 줄어들거나 사라졌으면 위반 (예: '챗지피티는' → '챗GPT는')
    2. 원문에선 자연스러운 끝(공백/문장끝/비한글)에 위치했는데
       교정문에선 그런 위치가 사라졌으면 위반 (expansion 케이스)
       예: '한기대' → '한기대학교'
    """
    if not custom_dict:
        return False
    for word in custom_dict:
        if not word:
            continue
        count_k = k.count(word)
        if count_k == 0:
            continue
        count_v = v.count(word)
        
        # 규칙 1: 단어 개수가 줄어들면 위반
        if count_v < count_k:
            return True
        
        # 규칙 2: 원문에선 단어가 자연스럽게 끝났는데 교정문에선 그렇지 않으면 위반
        # (한기대 → 한기대학교 같은 expansion 검출)
        if _has_natural_ending(k, word) and not _has_natural_ending(v, word):
            return True
        if _has_natural_start(k, word) and not _has_natural_start(v, word):
            return True
    return False


_HANGUL_RE = re.compile(r'[가-힣]')


def _has_natural_ending(text, word):
    """text 안에 word가 등장하면서 그 직후가 한글이 아닌(=단어가 자연스럽게 끝난) 경우가 있는지."""
    if not word or not text:
        return False
    idx = 0
    while True:
        pos = text.find(word, idx)
        if pos == -1:
            return False
        after_pos = pos + len(word)
        if after_pos >= len(text) or not _HANGUL_RE.match(text[after_pos]):
            return True
        idx = pos + 1


def _has_natural_start(text, word):
    """text 안에 word가 등장하면서 그 직전이 한글이 아닌(=단어가 자연스럽게 시작한) 경우가 있는지."""
    if not word or not text:
        return False
    idx = 0
    while True:
        pos = text.find(word, idx)
        if pos == -1:
            return False
        if pos == 0 or not _HANGUL_RE.match(text[pos - 1]):
            return True
        idx = pos + 1


# ==========================================
# OpenAI 교정 (PPT)
# ==========================================
def get_openai_corrections_by_slide(prs, api_key, is_paid_tier=True, custom_dict=None,
                                    progress_callback=None, model="gpt-4o"):
    """
    슬라이드를 하나씩 읽어가면서 문맥을 바탕으로 OpenAI 교정안을 확보합니다.
    model: 'gpt-4o' (정확) 또는 'gpt-4o-mini' (빠르고 저렴)
    """
    client = OpenAI(api_key=api_key)
    global_corrections = {}
    
    system_prompt = _build_system_prompt(custom_dict, doc_kind="파워포인트 슬라이드")
    total_slides = len(prs.slides)
    
    for i, slide in enumerate(prs.slides):
        # 그룹 도형까지 재귀 순회하여 모든 텍스트 수집
        slide_texts = []
        for shape in iter_shapes(slide.shapes):
            # 일반 텍스트 도형 (placeholder, autoshape 등 특수 도형도 안전하게 처리)
            t = _safe_shape_text(shape).strip()
            if t and len(t) > 1:
                slide_texts.append(t)
            # 표 셀
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            slide_texts.append(cell.text.strip())
                            
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_texts.append(notes_text)
                            
        full_text = "\n".join([t for t in slide_texts if len(t) > 1])
        
        if not full_text.strip():
            if progress_callback: progress_callback(i + 1, total_slides)
            continue
            
        user_prompt = f'=== 슬라이드 {i+1} 텍스트 ===\n{full_text}'

        success = False
        for attempt in range(5):
            try:
                response = client.chat.completions.create(
                    model=model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt}
                    ],
                    response_format={"type": "json_object"},
                    temperature=0.1
                )
                
                res_text = response.choices[0].message.content.strip()
                slide_dict = json.loads(res_text)
                
                for k, v in slide_dict.items():
                    k_str = str(k).strip()
                    v_str = str(v).strip()
                    
                    if not k_str or not v_str: continue
                    if k_str == v_str: continue
                    # 단일 부호·공백만 제외 (1글자짜리 정상 한글 교정은 살림)
                    if len(k_str) == 1 and k_str in {" ", ".", ",", "!", "?", "-", "_", "·", "/"}: continue
                    
                    if _is_custom_dict_violation(k_str, v_str, custom_dict):
                        print(f"   [사용자 사전] 슬라이드 {i+1} 교정 차단: '{k_str}' → '{v_str}'")
                        continue
                    
                    # 충돌 시 덮어쓰지 않음 (먼저 등록된 교정을 보존)
                    if k_str not in global_corrections:
                        global_corrections[k_str] = v_str
                
                success = True
                break
                
            except Exception as e:
                err_msg = str(e)
                if "rate limit" in err_msg.lower() or "429" in err_msg:
                    print(f"   [API 한도 초과] 5초 대기 후 슬라이드 {i+1} 재시도... ({attempt+1}/5)")
                    time.sleep(5) 
                else:
                    print(f"   [API 오류] 재시도 중... 사유: {e}")
                    time.sleep(2)
                    
        if success and not is_paid_tier:
            time.sleep(1)
            
        if progress_callback:
            progress_callback(i + 1, total_slides)
            
    return global_corrections


# ==========================================
# 교정 적용 (PPT)
# ==========================================
def apply_corrections_to_ppt(prs, corrections_dict):
    """
    교정 딕셔너리를 PPT 내부 텍스트에 적용하고, 변경된 부분을 핫핑크색으로 강조한다.
    그룹 도형까지 재귀적으로 처리한다.
    """
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    _apply_to_paragraph(paragraph, corrections_dict)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                _apply_to_paragraph(paragraph, corrections_dict)
                                
        if slide.has_notes_slide:
            for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                _apply_to_paragraph(paragraph, corrections_dict)


def _apply_to_paragraph(paragraph, corrections_dict):
    original_text = paragraph.text.strip()
    if not original_text:
        return
        
    corrected_text = original_text
    is_changed = False
    
    # 긴 키부터 치환해서 짧은 단어가 다른 단어를 망가뜨리는 것을 방지
    for old_txt, new_txt in _sorted_corrections(corrections_dict):
        if old_txt in corrected_text:
            corrected_text = corrected_text.replace(old_txt, new_txt)
            is_changed = True
            
    # 다중 띄어쓰기 정리
    spaced_fixed = re.sub(r' {2,}', ' ', corrected_text)
    if spaced_fixed != corrected_text:
        corrected_text = spaced_fixed
        is_changed = True
    
    # 구두점 정리 (이중 마침표 등 GPT 교정 부작용 처리)
    cleaned = _clean_punctuation(corrected_text)
    if cleaned != corrected_text:
        corrected_text = cleaned
        is_changed = True
        
    if not is_changed:
        return
        
    if paragraph.runs:
        font_ref = paragraph.runs[0].font
    else:
        font_ref = None
        
    paragraph.clear()
    
    tokens_orig = re.split(r'(\s+)', original_text)
    tokens_corr = re.split(r'(\s+)', corrected_text)
    
    matcher = difflib.SequenceMatcher(None, tokens_orig, tokens_corr)
    
    for tag, i1, i2, j1, j2 in matcher.get_opcodes():
        if tag == 'delete':
            continue
            
        chunk_text = "".join(tokens_corr[j1:j2])
        if not chunk_text:
            continue
            
        new_run = paragraph.add_run()
        new_run.text = chunk_text
        
        if font_ref:
            for attr in ['name', 'size', 'bold', 'italic', 'underline']:
                try: setattr(new_run.font, attr, getattr(font_ref, attr))
                except: pass
                
        if tag in ('replace', 'insert'):
            new_run.font.color.rgb = RGBColor(255, 105, 180)
        else:
            if font_ref and hasattr(font_ref, 'color') and hasattr(font_ref.color, 'rgb') and font_ref.color.rgb:
                try: new_run.font.color.rgb = font_ref.color.rgb
                except: pass
            elif font_ref and hasattr(font_ref, 'color') and hasattr(font_ref.color, 'theme_color') and font_ref.color.theme_color:
                try: new_run.font.color.theme_color = font_ref.color.theme_color
                except: pass


# ==========================================
# [PDF 전용 기능] PyMuPDF (fitz) 활용
# ==========================================

def extract_narrations_pdf(pdf_document):
    """
    PDF의 각 페이지 텍스트 블록을 스캔하여 화자 패턴을 분리한다.
    명시적 화자 라벨('교수:', '내레이션:' 등)이 있는 블록만 수집한다.
    """
    narrations = {"교수": [], "성우": [], "선생님": [], "기타": []}
    
    for i in range(len(pdf_document)):
        page = pdf_document[i]
        text_blocks = page.get_text("blocks")
        
        for block in text_blocks:
            text = block[4].strip()
            if not text:
                continue
            
            match = re.match(
                r'^\s*(교수|성우|선생님|내레이션|나레이션)님?\s*[:]\s*(.*)',
                text,
                flags=re.DOTALL
            )
            if match:
                label = match.group(1)
                narration_text = match.group(2).strip()
                if not narration_text:
                    continue
                speaker_found = _detect_speaker(label) or "기타"
                narrations[speaker_found].append(
                    f"[페이지 {i+1}] {label} :\n{narration_text}"
                )
            # 라벨 없는 블록은 수집하지 않음 (버전 정보, 화면 설명, 본문 등 노이즈 차단)
                    
    return narrations


def get_openai_corrections_by_page_pdf(pdf_document, api_key, is_paid_tier=True, custom_dict=None,
                                       progress_callback=None, model="gpt-4o"):
    """
    PDF 페이지별 텍스트를 추출해 OpenAI 교정안(JSON)을 받아옵니다.
    """
    client = OpenAI(api_key=api_key)
    global_corrections = {}
    
    system_prompt = _build_system_prompt(custom_dict, doc_kind="PDF 문서")
    total_pages = len(pdf_document)
    
    for i in range(total_pages):
        page = pdf_document[i]
        blocks = page.get_text("blocks")
        text_lines = []
        for b in blocks:
            if len(b) >= 7 and b[6] == 0:  # 0번이 텍스트 블록
                block_txt = b[4].strip()
                if len(block_txt) > 2 and not block_txt.isdigit():
                    text_lines.append(block_txt)
        
        full_text = "\n\n".join(text_lines)
        
        if not full_text:
            if progress_callback: progress_callback(i + 1, total_pages)
            continue
            
        user_prompt = f'=== 페이지 {i+1} 텍스트 ===\n{full_text}'

        success = False
        for attempt in range(5):
            try:
                response = client.chat.completions.create(
                    model=model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt}
                    ],
                    response_format={"type": "json_object"},
                    temperature=0.1
                )
                
                res_text = response.choices[0].message.content.strip()
                page_dict = json.loads(res_text)
                
                for k, v in page_dict.items():
                    k_str = str(k).strip()
                    v_str = str(v).strip()
                    
                    if not k_str or not v_str: continue
                    if k_str == v_str: continue
                    if len(k_str) == 1 and k_str in {" ", ".", ",", "!", "?", "-", "_", "·", "/"}: continue
                    
                    if _is_custom_dict_violation(k_str, v_str, custom_dict):
                        print(f"   [사용자 사전] 페이지 {i+1} 교정 차단: '{k_str}' → '{v_str}'")
                        continue
                    
                    if k_str not in global_corrections:
                        global_corrections[k_str] = v_str
                
                success = True
                break
                
            except Exception as e:
                err_msg = str(e)
                if "rate limit" in err_msg.lower() or "429" in err_msg:
                    time.sleep(5) 
                else:
                    time.sleep(2)
                    
        if success and not is_paid_tier:
            time.sleep(1)
            
        if progress_callback:
            progress_callback(i + 1, total_pages)
            
    return global_corrections


def apply_corrections_to_pdf(pdf_document, corrections_dict):
    """
    교정 딕셔너리를 바탕으로 PDF 원문에 핫핑크색 하이라이트 어노테이션을 그린다.
    긴 키부터 처리해서 짧은 키가 긴 키 안에 중복 매칭되는 것을 방지.
    """
    sorted_items = _sorted_corrections(corrections_dict)
    
    for page in pdf_document:
        for old_txt, new_txt in sorted_items:
            if not old_txt.strip(): continue
            
            text_instances = page.search_for(old_txt, quads=True)
            for inst in text_instances:
                highlight = page.add_highlight_annot(inst)
                highlight.set_colors(stroke=(1.0, 0.41, 0.70))
                
                info = highlight.info
                info["title"] = "AI 맞춤법 교정"
                info["content"] = f"제안: {new_txt}"
                highlight.set_info(info)
                
                highlight.update()
